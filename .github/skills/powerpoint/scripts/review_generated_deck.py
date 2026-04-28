"""Quality gate for generated stakeholder PowerPoint decks.

Evaluates generated decks against template consistency and stakeholder readability
heuristics. Emits JSON and Markdown review artifacts so poor-quality decks can be
blocked during generation.
"""

from __future__ import annotations

import argparse
import json
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
import re

from pptx import Presentation  # type: ignore[import-not-found]


FORBIDDEN_TOKENS = (
    "todo",
    "criterion 1",
    "criterion 2",
    "criterion 3",
    "replace this",
    "lorem ipsum",
)


@dataclass(frozen=True)
class CriterionResult:
    """One scored quality criterion."""

    name: str
    score: float
    rationale: str


def _slide_text(shape: Any) -> list[str]:
    """Extract visible text lines from one shape."""
    if not getattr(shape, "has_text_frame", False):
        return []
    lines: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)
    return lines


def _slide_lines(slide: Any) -> list[str]:
    """Collect all non-empty text lines from a slide."""
    lines: list[str] = []
    for shape in slide.shapes:
        lines.extend(_slide_text(shape))
    return lines


def _backup_index(presentation: Any) -> int:
    """Return BACKUP divider index, fallback to full deck length."""
    for index, slide in enumerate(presentation.slides):
        lines = _slide_lines(slide)
        if any(line.strip().upper() == "BACKUP" for line in lines):
            return index
    return len(presentation.slides)


def _ratio_score(successes: int, total: int) -> float:
    """Map pass ratio to a 1-5 score."""
    if total <= 0:
        return 1.0
    ratio = max(0.0, min(1.0, successes / total))
    return round(1.0 + (ratio * 4.0), 2)


def _score_template_alignment(deck: Any, template: Any, backup_idx: int) -> CriterionResult:
    """Score visual/structural alignment to template contract."""
    checks = 0
    passed = 0

    checks += 1
    if deck.slide_width == template.slide_width and deck.slide_height == template.slide_height:
        passed += 1

    checks += 1
    if backup_idx < len(deck.slides):
        passed += 1

    checks += 1
    if backup_idx >= 5:
        passed += 1

    checks += 1
    if len(deck.slides) >= len(template.slides):
        passed += 1

    score = _ratio_score(passed, checks)
    return CriterionResult(
        name="Template Compliance",
        score=score,
        rationale=(
            f"{passed}/{checks} structural checks passed (size, backup divider, "
            "minimum front section, template footprint)."
        ),
    )


def _score_story_flow(deck: Any, backup_idx: int) -> CriterionResult:
    """Score whether front section follows a clear narrative flow."""
    front_count = min(backup_idx, len(deck.slides))
    checks = 0
    passed = 0

    checks += 1
    if front_count >= 5:
        passed += 1

    for idx in range(min(front_count, 5)):
        checks += 1
        lines = _slide_lines(deck.slides[idx])
        if lines and len(" ".join(lines[:3])) >= 12:
            passed += 1

    score = _ratio_score(passed, checks)
    return CriterionResult(
        name="Narrative Structure",
        score=score,
        rationale=(
            f"{passed}/{checks} checks passed for title/agenda/chapter/content/closing "
            "readability in the front section."
        ),
    )


def _is_bullet_line(line: str) -> bool:
    """Heuristic bullet detection for extracted text lines."""
    stripped = line.strip()
    return stripped.startswith(("-", "*", "•"))


def _score_text_density(deck: Any, backup_idx: int) -> CriterionResult:
    """Score text density and bullet hygiene for stakeholder slides."""
    front_count = min(backup_idx, len(deck.slides))
    checks = 0
    passed = 0

    for idx in range(front_count):
        lines = _slide_lines(deck.slides[idx])
        if not lines:
            continue

        body_lines = lines[1:] if len(lines) > 1 else []
        bullet_count = sum(1 for line in body_lines if _is_bullet_line(line))
        total_chars = sum(len(line) for line in body_lines)

        checks += 1
        if bullet_count <= 3:
            passed += 1

        checks += 1
        if total_chars <= 420:
            passed += 1

    score = _ratio_score(passed, checks)
    return CriterionResult(
        name="Text Density",
        score=score,
        rationale=(
            f"{passed}/{checks} checks passed for concise body text and <=3 bullet guideline."
        ),
    )


def _score_placeholder_hygiene(deck: Any, backup_idx: int) -> CriterionResult:
    """Score placeholder/noise token cleanliness."""
    front_count = min(backup_idx, len(deck.slides))
    checks = 0
    passed = 0

    for idx in range(front_count):
        lines = _slide_lines(deck.slides[idx])
        for line in lines:
            checks += 1
            lowered = line.lower()
            if not any(token in lowered for token in FORBIDDEN_TOKENS):
                passed += 1

    score = _ratio_score(passed, checks)
    return CriterionResult(
        name="Content Hygiene",
        score=score,
        rationale=f"{passed}/{checks} text lines passed placeholder/noise-token validation.",
    )


def _score_user_clarity(deck: Any, backup_idx: int) -> CriterionResult:
    """Score from a simple non-technical user perspective."""
    front_count = min(backup_idx, len(deck.slides))
    checks = 0
    passed = 0

    for idx in range(front_count):
        lines = _slide_lines(deck.slides[idx])
        if not lines:
            continue

        title = lines[0].strip()
        checks += 1
        if 12 <= len(title) <= 90:
            passed += 1

        checks += 1
        has_action_or_outcome = any(
            marker in " ".join(lines).lower()
            for marker in (
                "next",
                "decision",
                "risk",
                "goal",
                "owner",
                "status",
                "scope",
            )
        )
        if has_action_or_outcome:
            passed += 1

    score = _ratio_score(passed, checks)
    return CriterionResult(
        name="User Lens",
        score=score,
        rationale=(
            f"{passed}/{checks} clarity checks passed for headline quality and decision-oriented wording."
        ),
    )


def _topic_tokens_from_source(source_path: Path) -> list[str]:
    """Extract a small set of topic tokens from source markdown/text."""
    if not source_path.exists() or not source_path.is_file():
        return []
    text = source_path.read_text(encoding="utf-8", errors="replace").lower()
    token_set: set[str] = set()
    if "traveling salesman" in text or "tsp" in text:
        token_set.update({"tsp", "traveling", "salesman", "quantum", "classical"})
    for token in ("notebook", "qaoa", "vqe", "qiskit", "hypergraph", "delivery", "scope"):
        if token in text:
            token_set.add(token)
    if token_set:
        return sorted(token_set)

    words = [word for word in re.findall(r"[a-z0-9]{4,}", text) if len(word) >= 4]
    # Keep deterministic order with simple frequency ranking.
    counts: dict[str, int] = {}
    for word in words:
        counts[word] = counts.get(word, 0) + 1
    ranked = sorted(counts.items(), key=lambda item: (-item[1], item[0]))
    return [word for word, _ in ranked[:5]]


def _score_source_relevance(deck: Any, backup_idx: int, source_path: Path) -> CriterionResult:
    """Score whether front-slide language reflects the source topic."""
    topics = _topic_tokens_from_source(source_path)
    if not topics:
        return CriterionResult(
            name="Source Relevance",
            score=3.0,
            rationale="No stable source tokens detected; relevance scoring skipped.",
        )

    front_count = min(backup_idx, len(deck.slides))
    front_text = "\n".join(
        " ".join(_slide_lines(deck.slides[index])).lower()
        for index in range(front_count)
    )
    matched = sum(1 for token in topics if token in front_text)
    score = _ratio_score(matched, max(1, min(5, len(topics))))
    return CriterionResult(
        name="Source Relevance",
        score=score,
        rationale=f"{matched}/{len(topics)} source topic tokens detected in front-section slide text.",
    )


def _score_screenshot_coverage(deck: Any, backup_idx: int, screenshots_dir: Path) -> CriterionResult:
    """Score whether screenshot rendering covers the complete front section."""
    if not screenshots_dir.exists() or not screenshots_dir.is_dir():
        return CriterionResult(
            name="Screenshot Coverage",
            score=1.0,
            rationale="Screenshot directory missing; visual review evidence unavailable.",
        )

    files = [
        path
        for path in screenshots_dir.iterdir()
        if path.is_file() and path.suffix.lower() in {".png", ".jpg", ".jpeg", ".webp"}
    ]
    front_count = min(backup_idx, len(deck.slides))
    score = _ratio_score(min(len(files), front_count), max(1, front_count))
    return CriterionResult(
        name="Screenshot Coverage",
        score=score,
        rationale=f"{len(files)}/{front_count} expected front slides have rendered screenshot files.",
    )


def review_deck(
    deck_path: Path,
    template_path: Path,
    source_path: Path | None = None,
    screenshots_dir: Path | None = None,
) -> dict[str, Any]:
    """Run full review and return structured payload."""
    deck = Presentation(str(deck_path))
    template = Presentation(str(template_path))
    backup_idx = _backup_index(deck)

    criteria = [
        _score_template_alignment(deck, template, backup_idx),
        _score_story_flow(deck, backup_idx),
        _score_text_density(deck, backup_idx),
        _score_placeholder_hygiene(deck, backup_idx),
        _score_user_clarity(deck, backup_idx),
    ]

    if source_path is not None:
        criteria.append(_score_source_relevance(deck, backup_idx, source_path))
    if screenshots_dir is not None:
        criteria.append(_score_screenshot_coverage(deck, backup_idx, screenshots_dir))

    composite = round(sum(item.score for item in criteria) / len(criteria), 2)
    if composite >= 4.0:
        recommendation = "proceed"
    elif composite >= 2.0:
        recommendation = "revise"
    else:
        recommendation = "redesign"

    return {
        "deck": str(deck_path),
        "template": str(template_path),
        "generated_at": datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"),
        "slides_total": len(deck.slides),
        "slides_front": backup_idx,
        "backup_present": backup_idx < len(deck.slides),
        "composite_score": composite,
        "recommendation": recommendation,
        "criteria": [
            {"name": item.name, "score": item.score, "rationale": item.rationale}
            for item in criteria
        ],
    }


def _review_markdown(payload: dict[str, Any]) -> str:
    """Render reviewer payload as stakeholder-friendly markdown."""
    lines = [
        "# PowerPoint Quality Review",
        "",
        f"- generated_at: {payload['generated_at']}",
        f"- deck: {payload['deck']}",
        f"- template: {payload['template']}",
        f"- slides_total: {payload['slides_total']}",
        f"- slides_front: {payload['slides_front']}",
        f"- backup_present: {str(payload['backup_present']).lower()}",
        f"- composite_score: {payload['composite_score']}",
        f"- recommendation: {payload['recommendation']}",
        "",
        "## Criteria",
        "",
        "| Criterion | Score (1-5) | Rationale |",
        "|---|---:|---|",
    ]
    for item in payload["criteria"]:
        lines.append(f"| {item['name']} | {item['score']} | {item['rationale']} |")
    return "\n".join(lines) + "\n"


def parse_args() -> argparse.Namespace:
    """Parse CLI arguments for deck review."""
    parser = argparse.ArgumentParser(description="Review generated PowerPoint quality")
    parser.add_argument("--deck", required=True, help="Generated deck path")
    parser.add_argument("--template", required=True, help="Template path used for generation")
    parser.add_argument("--source", help="Original source path used for generation")
    parser.add_argument("--screenshots-dir", help="Rendered screenshot directory for the generated deck")
    parser.add_argument("--min-score", type=float, default=4.0, help="Minimum composite score")
    parser.add_argument("--report-json", required=True, help="Output JSON report path")
    parser.add_argument("--report-md", required=True, help="Output Markdown report path")
    parser.add_argument(
        "--strict",
        action="store_true",
        help="Exit non-zero if quality recommendation is not proceed or min-score is not met",
    )
    return parser.parse_args()


def main() -> int:
    """CLI entrypoint."""
    args = parse_args()
    deck_path = Path(args.deck).expanduser().resolve()
    template_path = Path(args.template).expanduser().resolve()
    report_json = Path(args.report_json).expanduser().resolve()
    report_md = Path(args.report_md).expanduser().resolve()
    source_path = Path(args.source).expanduser().resolve() if args.source else None
    screenshots_dir = (
        Path(args.screenshots_dir).expanduser().resolve()
        if args.screenshots_dir
        else None
    )

    if not deck_path.exists():
        raise SystemExit(f"deck not found: {deck_path}")
    if not template_path.exists():
        raise SystemExit(f"template not found: {template_path}")

    payload = review_deck(
        deck_path,
        template_path,
        source_path=source_path,
        screenshots_dir=screenshots_dir,
    )
    report_json.parent.mkdir(parents=True, exist_ok=True)
    report_md.parent.mkdir(parents=True, exist_ok=True)

    report_json.write_text(json.dumps(payload, indent=2), encoding="utf-8")
    report_md.write_text(_review_markdown(payload), encoding="utf-8")

    output = {
        "status": "ok",
        "quality_review": {
            "json": str(report_json),
            "markdown": str(report_md),
            "composite_score": payload["composite_score"],
            "recommendation": payload["recommendation"],
            "min_score": args.min_score,
            "strict": bool(args.strict),
        },
    }
    print(json.dumps(output))

    if args.strict:
        if payload["composite_score"] < args.min_score or payload["recommendation"] != "proceed":
            return 3
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
