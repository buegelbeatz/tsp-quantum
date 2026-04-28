"""Build a presentation from source content using a layer template."""

from __future__ import annotations

import argparse
import copy
import json
import re
from pathlib import Path
import sys
import tempfile
from typing import Any

if str(Path(__file__).resolve().parent) not in sys.path:
    sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation  # type: ignore[import-not-found]
from pptx.dml.color import RGBColor  # type: ignore[import-not-found]
from pptx.enum.text import PP_ALIGN  # type: ignore[import-not-found]
from pptx.util import Inches, Pt  # type: ignore[import-not-found]

from powerpoint_seed import source_slug  # type: ignore[import-not-found]
from template_factory import ensure_template  # type: ignore[import-not-found]

try:
    import cairosvg as _cairosvg  # type: ignore[import-not-found]
except ImportError:  # pragma: no cover - optional dependency fallback
    _cairosvg = None

CAIROSVG = _cairosvg


def _collect_chunks(source: Path) -> list[tuple[str, str]]:
    if source.is_file():
        text = source.read_text(encoding="utf-8", errors="replace")
        if source.suffix.lower() == ".md":
            sections = _sections_from_markdown(text, source.stem)
            if sections:
                return sections
        return [(source.stem, _clean_text("\n".join(text.splitlines()[:18])))]

    chunks: list[tuple[str, str]] = []
    for path in sorted(source.rglob("*")):
        if (
            len(chunks) >= 8
            or not path.is_file()
            or path.suffix.lower() not in {".md", ".txt"}
        ):
            continue
        text = path.read_text(encoding="utf-8", errors="replace")
        if path.suffix.lower() == ".md":
            for title, body in _sections_from_markdown(text, path.stem):
                chunks.append((title, body))
                if len(chunks) >= 8:
                    break
            continue
        body = _clean_text("\n".join(text.splitlines()[:14]))
        chunks.append((path.stem, body or "Content extracted from source text."))
    return chunks or [(source.name, "No readable markdown or text files found.")]


def _sections_from_markdown(markdown_text: str, fallback_title: str) -> list[tuple[str, str]]:
    """Convert markdown headings into structured slide sections.

    Keeps heading and body separated so slides do not collapse into unstructured
    fragments.
    """
    sections: list[tuple[str, str]] = []
    heading_pattern = re.compile(r"^(#{1,3})\s+(.+?)\s*$")
    current_title = fallback_title
    current_body: list[str] = []

    def flush() -> None:
        cleaned = _clean_text("\n".join(current_body))
        if cleaned:
            sections.append((current_title, cleaned))

    for raw_line in markdown_text.splitlines():
        match = heading_pattern.match(raw_line.strip())
        if match:
            if current_body:
                flush()
                current_body = []
            heading_text = match.group(2).strip()
            if heading_text:
                current_title = heading_text
            continue
        current_body.append(raw_line)

    if current_body:
        flush()

    if not sections:
        fallback = _clean_text(markdown_text)
        if fallback:
            return [(fallback_title, fallback)]
    return sections[:8]


def _extract_section(markdown_text: str, heading: str) -> str:
    """Extract markdown section body by heading name."""
    lines = markdown_text.splitlines()
    in_section = False
    collected: list[str] = []
    expected = heading.strip().lower()
    for line in lines:
        if line.startswith("## "):
            current = line[3:].strip().lower()
            if in_section:
                break
            if current == expected:
                in_section = True
                continue
        if in_section:
            collected.append(line)
    return "\n".join(collected).strip()


def _bullet_lines(section_text: str) -> list[str]:
    """Extract bullet items from markdown section text."""
    items: list[str] = []
    for line in section_text.splitlines():
        stripped = line.strip()
        if stripped.startswith(("- ", "* ")):
            text = stripped[2:].strip()
            if text:
                items.append(text)
    return items


def _clean_text(text: str) -> str:
    """Remove placeholder lines and normalize markdown text for slides."""
    blocked_tokens = (
        "todo",
        "criterion 1",
        "criterion 2",
        "criterion 3",
        "define the explicitly included outcomes",
        "define explicit exclusions",
        "replace this",
    )
    lines: list[str] = []
    for raw in text.splitlines():
        line = raw.strip()
        if not line:
            continue
        normalized = line.lower()
        if any(token in normalized for token in blocked_tokens):
            continue
        line = re.sub(r"^[#>-]+\s*", "", line).strip()
        line = re.sub(r"^(primary finding:|secondary finding:|core problem:|supporting evidence:|address:|focus on:)\s*", "", line, flags=re.IGNORECASE)
        line = line.replace("**", "").replace("`", "").strip()
        if "|" in line and line.count("|") >= 2:
            continue
        if line:
            lines.append(line)
    return "\n".join(lines)


def _slide_palette(slide_index: int) -> tuple[RGBColor, RGBColor, RGBColor, RGBColor]:
    """Return title, eyebrow, body, divider colors for the template slide slot."""
    if slide_index in {3, 6}:
        return (
            RGBColor(0x1F, 0x4E, 0x79),
            RGBColor(0x16, 0xE0, 0xBD),
            RGBColor(0x07, 0x14, 0x26),
            RGBColor(0xCF, 0xD8, 0xE3),
        )
    return (
        RGBColor(0xFF, 0xFF, 0xFF),
        RGBColor(0x16, 0xE0, 0xBD),
        RGBColor(0xE5, 0xE7, 0xEB),
        RGBColor(0x3B, 0xD8, 0xC5),
    )


def _set_text_style(
    paragraph: Any,
    *,
    font_name: str | None,
    size: int,
    color: RGBColor,
    bold: bool = False,
) -> None:
    """Apply a consistent text style to all runs in a paragraph."""
    for run in paragraph.runs:
        if font_name:
            run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _fit_body_lines(lines: list[str], max_lines: int = 10, max_chars: int = 98) -> list[str]:
    """Trim body lines to reduce text overflow in fixed-size text boxes."""
    fitted = [line[:max_chars].rstrip() for line in lines if line.strip()]
    if len(fitted) <= max_lines:
        return fitted
    clipped = fitted[: max_lines - 1]
    clipped.append("Additional details are documented in the source artifact")
    return clipped


def _placeholder_index(shape: Any) -> int | None:
    """Return placeholder index for placeholder shapes, otherwise None."""
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        return int(shape.placeholder_format.idx)
    except (AttributeError, TypeError, ValueError):
        return None


def _render_body(text_frame: Any, body: str, body_color: RGBColor) -> None:
    """Render body text with bullet-aware formatting."""
    text_frame.clear()
    lines = [line.strip() for line in body.splitlines() if line.strip()]
    if not lines:
        lines = ["Content is being refined from the canonical stage document."]
    lines = _fit_body_lines(lines)
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        bullet_line = line.startswith("-")
        paragraph.text = line.lstrip("- ").strip() if bullet_line else line
        paragraph.level = 0
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(6)
        if bullet_line:
            paragraph.bullet = True
        _set_text_style(paragraph, font_name=None, size=16, color=body_color)


def _populate_template_slide(
    slide: Any,
    slide_index: int,
    title: str,
    body: str,
    visual_asset: Path | None = None,
) -> None:
    """Populate one template slide.

    Strategy: fill the template slide's existing text frames in-place, preserving
    the template's visual design (backgrounds, accents, imagery).  New textboxes
    are only added as a last-resort fallback when the slide has no writable text
    frames at all.
    """
    title_color, _eyebrow_color, body_color, _divider_color = _slide_palette(slide_index)

    title_placed = False
    body_placed = False

    # --- Pass 1: fill standard PowerPoint placeholders (idx 0 = title, idx 1 = body) ---
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        ph_idx = _placeholder_index(shape)
        if ph_idx is None:
            continue
        if ph_idx == 0 and not title_placed:
            shape.text_frame.clear()
            shape.text_frame.text = title
            _set_text_style(
                shape.text_frame.paragraphs[0],
                font_name=None,
                size=26,
                color=title_color,
                bold=True,
            )
            title_placed = True
        elif ph_idx == 1 and not body_placed:
            _render_body(shape.text_frame, body, body_color)
            body_placed = True
        if title_placed and body_placed:
            break

    # --- Pass 2: fill non-placeholder text frames when placeholders were absent ---
    if not title_placed or not body_placed:
        text_shapes = [
            s for s in slide.shapes
            if getattr(s, "has_text_frame", False)
            and _placeholder_index(s) is None
        ]
        # Heuristic: widest tall-ish shape → title; tallest/largest shape → body
        text_shapes_sorted = sorted(
            text_shapes,
            key=lambda s: (s.width * s.height),
            reverse=True,
        )
        for shape in text_shapes_sorted:
            if not title_placed:
                shape.text_frame.clear()
                shape.text_frame.text = title
                _set_text_style(
                    shape.text_frame.paragraphs[0],
                    font_name=None,
                    size=26,
                    color=title_color,
                    bold=True,
                )
                title_placed = True
            elif not body_placed:
                _render_body(shape.text_frame, body, body_color)
                body_placed = True
                break

    # --- Fallback: add textboxes only when the template had no usable text frames ---
    if not title_placed:
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.78), Inches(11.8), Inches(0.9))
        title_box.text_frame.text = title
        _set_text_style(
            title_box.text_frame.paragraphs[0],
            font_name=None,
            size=26,
            color=title_color,
            bold=True,
        )
    if not body_placed:
        body_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(11.2), Inches(4.9))
        _render_body(body_box.text_frame, body, body_color)

    if visual_asset is not None:
        visual_left = Inches(8.45)
        visual_top = Inches(4.15)
        visual_width = Inches(3.5)
        visual_height = Inches(2.2)
        text_overlaps = False
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            sx = int(getattr(shape, "left", 0))
            sy = int(getattr(shape, "top", 0))
            sw = int(getattr(shape, "width", 0))
            sh = int(getattr(shape, "height", 0))
            if sw <= 0 or sh <= 0:
                continue
            if (
                sx < visual_left + visual_width
                and sx + sw > visual_left
                and sy < visual_top + visual_height
                and sy + sh > visual_top
            ):
                text_overlaps = True
                break

        if text_overlaps:
            return
        try:
            slide.shapes.add_picture(
                str(visual_asset),
                visual_left,
                visual_top,
                width=visual_width,
                height=visual_height,
            )
        except (ValueError, OSError, TypeError, AttributeError):
            # Best-effort only: keep slide generation robust when asset format is unsupported.
            pass


def _add_content_slide(prs: Any, title: str, body: str) -> None:
    """Append a readable content slide with clear heading/body separation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    bg.line.color.rgb = RGBColor(0xF8, 0xFA, 0xFC)

    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.7), Inches(11.8), Inches(0.8)
    )
    title_box.text_frame.text = title
    title_run = title_box.text_frame.paragraphs[0].runs[0]
    title_run.font.size = Pt(28)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    body_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.8), Inches(11.6), Inches(5.3)
    )
    _render_body(body_box.text_frame, body, RGBColor(0x07, 0x14, 0x26))


def _duplicate_slide(prs: Any, source_index: int) -> int:
    """Duplicate a slide and return the index of the appended copy.

    The copied slide preserves the original template visuals while allowing
    front-area population on the duplicate only.
    """
    source_slide = prs.slides[source_index]
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)

    rel_map: dict[str, str] = {}
    rel_namespace = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    embed_attr = f"{{{rel_namespace}}}embed"
    link_attr = f"{{{rel_namespace}}}link"
    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        if rel.is_external:
            new_rid = new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rid = new_slide.part.rels.get_or_add(rel.reltype, rel._target)
        rel_map[rel.rId] = new_rid

    for shape in source_slide.shapes:
        new_element = copy.deepcopy(shape.element)
        for element in new_element.iter():
            old_embed = element.get(embed_attr)
            if old_embed in rel_map:
                element.set(embed_attr, rel_map[old_embed])
            old_link = element.get(link_attr)
            if old_link in rel_map:
                element.set(link_attr, rel_map[old_link])
        new_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")

    return len(prs.slides) - 1


def _discover_visual_assets(repo_root: Path, source: Path) -> list[Path]:
    """Discover a small set of visual assets to enrich stakeholder slides."""
    def _topic_tokens_from_source(path: Path) -> set[str]:
        tokens: set[str] = set()
        name_tokens = re.findall(r"[a-z0-9]+", path.stem.lower())
        tokens.update(token for token in name_tokens if len(token) >= 3)
        if path.exists() and path.is_file() and path.suffix.lower() in {".md", ".txt"}:
            text = path.read_text(encoding="utf-8", errors="replace").lower()
            if "traveling salesman" in text or "tsp" in text:
                tokens.update({"tsp", "traveling", "salesman", "quantum", "classical"})
            for token in ("notebook", "qaoa", "vqe", "qiskit", "hypergraph"):
                if token in text:
                    tokens.add(token)
        return tokens

    def _path_matches_topic(path: Path, topic_tokens: set[str]) -> bool:
        if not topic_tokens:
            return False
        lowered = path.as_posix().lower()
        stem_tokens = set(re.findall(r"[a-z0-9]+", path.stem.lower()))
        return any(token in lowered for token in topic_tokens) or bool(stem_tokens & topic_tokens)

    def _stage_from_markdown(path: Path) -> str:
        if not path.exists() or path.suffix.lower() != ".md":
            return "project"
        text = path.read_text(encoding="utf-8", errors="replace")
        match = re.search(r'^stage:\s*"?([a-z\-]+)"?\s*$', text, flags=re.MULTILINE)
        if match:
            return match.group(1).strip().lower()
        return "project"

    stage_name = _stage_from_markdown(source)
    topic_tokens = _topic_tokens_from_source(source)
    roots = [
        source.parent / "assets",
        repo_root / ".digital-artifacts" / "40-stage" / "assets",
        repo_root / ".digital-artifacts" / "50-planning" / stage_name / "assets",
        repo_root / ".digital-artifacts" / "10-data",
        repo_root / ".digital-artifacts" / "60-review",
        repo_root / "docs" / "wiki" / "assets" / "visualizations",
        repo_root / "docs" / "images" / "mermaid",
        repo_root / "docs" / "wiki" / "assets",
    ]
    candidates: list[Path] = []
    allowed = {".png", ".jpg", ".jpeg", ".webp", ".svg"}
    for root in roots:
        if not root.exists() or not root.is_dir():
            continue
        for path in sorted(root.rglob("*")):
            if path.is_file() and path.suffix.lower() in allowed:
                lowered = path.as_posix().lower()
                if "docs/ux/scribbles" in lowered and not _path_matches_topic(path, topic_tokens):
                    continue
                if "digital-generic-team" in path.name.lower() and stage_name not in lowered:
                    continue
                if "portraits" in lowered:
                    continue
                if "powerpoint" in lowered and "quality-review" in lowered:
                    continue
                if "10-data" in lowered and not any(token in lowered for token in ("plot", "chart", "figure", "visual", "map", "graph", "output", "result", "assets")):
                    continue
                candidates.append(path)
                if len(candidates) >= 8:
                    return candidates
    return candidates


def _prepare_visual_asset_for_ppt(asset: Path, temp_dir: Path) -> Path | None:
    """Convert unsupported image formats to a PPT-friendly file when needed."""
    if asset.suffix.lower() != ".svg":
        return asset
    if CAIROSVG is None:
        return None

    converted = temp_dir / f"{asset.stem}-{abs(hash(str(asset)))}.png"
    try:
        CAIROSVG.svg2png(url=str(asset), write_to=str(converted))
    except (OSError, ValueError, TypeError, AttributeError):
        return None
    return converted if converted.exists() else None


def _stage_sections(source: Path, layer: str) -> list[tuple[str, str]]:
    """Build a concise stakeholder section list from canonical stage markdown."""
    stage_text = source.read_text(encoding="utf-8", errors="replace")
    title = ""
    for line in stage_text.splitlines():
        if line.startswith("# "):
            title = line[2:].strip()
            break
    stage_title = title or f"{layer.title()} Project"
    if len(stage_title) < 12:
        stage_title = f"{stage_title} Stage Status"

    vision = _clean_text(_extract_section(stage_text, "Vision"))
    goals = _bullet_lines(_clean_text(_extract_section(stage_text, "Goals")))
    constraints = _bullet_lines(_clean_text(_extract_section(stage_text, "Constraints")))
    readiness = _bullet_lines(
        _clean_text(_extract_section(stage_text, "Specification Readiness Summary"))
    )
    open_questions = _bullet_lines(_clean_text(_extract_section(stage_text, "Open Questions")))
    scope_boundaries = _clean_text(_extract_section(stage_text, "Scope Boundaries"))
    current_context = _clean_text(_extract_section(stage_text, "Current Context"))
    summary_lines = [line for line in vision.splitlines() if line.strip()][:3]
    if not summary_lines:
        summary_lines = ["This project stage consolidates validated scope into delivery-ready work."]

    def _is_noise(line: str) -> bool:
        lowered = line.lower()
        return any(
            marker in lowered
            for marker in (
                "task audit",
                "error occurred",
                "/stages-action",
                "rerun",
                "session id",
                "prompt-invoke",
                "next step is to check the command output",
            )
        )

    context_lines = [line for line in current_context.splitlines() if line.strip() and not _is_noise(line)][:5]
    scope_lines = [line for line in scope_boundaries.splitlines() if line.strip() and not _is_noise(line)][:6]
    goal_lines = [line for line in goals if line.strip() and not _is_noise(line)][:6]
    constraint_lines = [line for line in constraints if line.strip() and not _is_noise(line)][:6]
    question_lines = [line for line in open_questions if line.strip() and not _is_noise(line)][:4]

    topic_hint = ""
    for line in context_lines:
        lowered = line.lower()
        if "traveling salesman" in lowered or "tsp" in lowered:
            topic_hint = line
            break

    opening_title = topic_hint or stage_title
    opening_body = (
        "Project stage briefing\n- Purpose: decision-ready summary for scope, delivery, and review status.\n- Focus: transform validated context into executable delivery work."
    )
    if topic_hint:
        opening_body = (
            "Decision-ready project briefing\n"
            "- Focus: compare classical and quantum-inspired TSP approaches with transparent trade-offs.\n"
            "- Outcome: align implementation scope, evidence, and next decisions for delivery."
        )

    return [
        (
            opening_title,
            opening_body,
        ),
        (
            "Executive Summary",
            "\n".join(f"- {line}" for line in summary_lines),
        ),
        (
            "Current Context",
            "\n".join([f"- {item}" for item in context_lines])
            if context_lines
            else "- Current context is being consolidated from the stage document.",
        ),
        (
            "Scope and Delivery Focus",
            "\n".join(
                [
                    "Goals:",
                    *[f"- {item}" for item in goal_lines],
                    "",
                    "Scope boundaries:",
                    *[f"- {item}" for item in scope_lines],
                ]
            )
            if goal_lines or scope_lines
            else "- Scope and goals are being refined in the canonical stage document.",
        ),
        (
            "Constraints, Risks, and Next Step",
            "\n".join(
                [
                    "Constraints:",
                    *[f"- {item}" for item in constraint_lines],
                    "",
                    "Readiness summary:",
                    *[f"- {item}" for item in readiness[:4]],
                    "",
                    "Open questions:",
                    *([f"- {item}" for item in question_lines] or ["- none"]),
                ]
            ),
        ),
    ]


def _clear_slide_text(slide: Any) -> None:
    """Clear all text from existing text frames on a slide."""
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            shape.text_frame.clear()


def _remove_slide(prs: Any, index: int) -> None:
    """Remove one slide by index using python-pptx internals."""
    slides = getattr(prs.slides, "_sldIdLst")
    slide_id = slides[index]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    slides.remove(slide_id)


def _reorder_slides(prs: Any, new_order: list[int]) -> None:
    """Reorder all slides to the given index sequence.

    Args:
        prs: The python-pptx Presentation object.
        new_order: A permutation of ``range(len(prs.slides))`` that defines
            the desired slide order.  Every current index must appear exactly
            once.
    """
    xml_slides = prs.slides._sldIdLst  # noqa: SLF001
    original = list(xml_slides)
    for elem in original:
        xml_slides.remove(elem)
    for idx in new_order:
        xml_slides.append(original[idx])


_BACKUP_DARK = RGBColor(0x07, 0x14, 0x26)
_BACKUP_TEAL = RGBColor(0x16, 0xE0, 0xBD)
_BACKUP_LIGHT = RGBColor(0xE5, 0xE7, 0xEB)


def _add_backup_divider_slide(prs: Any) -> None:
    """Append a BACKUP divider slide that separates content from template sources.

    The slide uses a dark navy background with a teal "BACKUP" heading and a
    light-grey sub-caption so that reviewers can instantly identify the
    boundary between published content and template reference slides.

    Args:
        prs: The python-pptx Presentation object to which the slide is appended.
    """
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = _BACKUP_DARK
    bg.line.color.rgb = _BACKUP_DARK

    heading = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.4)
    )
    heading.text_frame.text = "BACKUP"
    run = heading.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = _BACKUP_TEAL
    run.font.name = "Aptos Display"

    caption = slide.shapes.add_textbox(
        Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.7)
    )
    caption.text_frame.text = "Template source slides below — for reference only"
    cap_run = caption.text_frame.paragraphs[0].runs[0]
    cap_run.font.size = Pt(16)
    cap_run.font.color.rgb = _BACKUP_LIGHT
    cap_run.font.name = "Aptos"


def _build_template_flow_sections(
    sections: list[tuple[str, str]],
) -> tuple[tuple[str, str], tuple[str, str], tuple[str, str], list[tuple[str, str]]]:
    """Map arbitrary source sections to template roles (title/agenda/chapter/content)."""
    if not sections:
        fallback = ("Project Briefing", "Decision-ready stage summary.")
        agenda = ("Agenda", "- Context and Objectives\n- Solution Overview\n- Next Steps and Decisions")
        chapter = ("Context and Objectives", "Current status, key assumptions, and scope boundaries.")
        return fallback, agenda, chapter, []

    title_section = sections[0]
    agenda_candidates = [title for title, _ in sections[1:4] if title.strip()]
    if not agenda_candidates:
        agenda_candidates = [
            "Context and Objectives",
            "Solution Overview",
            "Next Steps and Decisions",
        ]
    while len(agenda_candidates) < 3:
        agenda_candidates.append("Next Steps and Decisions")
    agenda_section = (
        "Agenda and Scope Decisions",
        "\n".join(
            f"{index}. {item}" for index, item in enumerate(agenda_candidates[:3], start=1)
        ),
    )

    chapter_source = sections[1] if len(sections) > 1 else sections[0]
    chapter_lines = [line.strip() for line in chapter_source[1].splitlines() if line.strip()][:3]
    chapter_body = "\n".join(f"- {line.lstrip('- ').strip()}" for line in chapter_lines) if chapter_lines else "- Key context and scope summary"
    chapter_section = (chapter_source[0], chapter_body)

    content_sections = sections[2:] if len(sections) > 2 else sections[1:]
    return title_section, agenda_section, chapter_section, content_sections


def _apply_sections_to_template(
    prs: Any, sections: list[tuple[str, str]], visual_assets: list[Path] | None = None
) -> None:
    """Populate selected template slides and move unused ones behind a BACKUP divider.

    Content sections are matched to preferred template slide positions.  When
    more sections than template slots exist, extra content slides are appended.
    Unused template slides (e.g. the Team or Schema reference slide) are NOT
    deleted; instead a styled BACKUP divider slide is inserted after all content
    slides so that the source material remains accessible for reviewers.

    Args:
        prs: python-pptx Presentation already loaded with the layer template.
        sections: Ordered list of ``(title, body)`` pairs to place in the deck.
        visual_assets: Optional ordered list of visual asset paths, aligned with
            *sections*.  Each asset is scaled and composited onto the matching
            slide when provided.
    """
    if not sections:
        return

    title_section, agenda_section, chapter_section, content_sections = _build_template_flow_sections(sections)

    template_count = len(prs.slides)
    template_indices = list(range(template_count))
    team_template_index = 4 if template_count > 4 else None
    thank_you_template_index = 5 if template_count > 5 else None
    content_template_index = 3 if template_count > 3 else None

    assets = visual_assets or []
    front_indices: list[int] = []

    if template_count > 0:
        copied_title = _duplicate_slide(prs, 0)
        _clear_slide_text(prs.slides[copied_title])
        _populate_template_slide(
            prs.slides[copied_title],
            0,
            title_section[0],
            title_section[1],
            None,
        )
        front_indices.append(copied_title)

    if template_count > 1:
        copied_agenda = _duplicate_slide(prs, 1)
        _clear_slide_text(prs.slides[copied_agenda])
        _populate_template_slide(
            prs.slides[copied_agenda],
            1,
            agenda_section[0],
            agenda_section[1],
            None,
        )
        front_indices.append(copied_agenda)

    if template_count > 2:
        copied_chapter = _duplicate_slide(prs, 2)
        _clear_slide_text(prs.slides[copied_chapter])
        _populate_template_slide(
            prs.slides[copied_chapter],
            2,
            chapter_section[0],
            chapter_section[1],
            None,
        )
        front_indices.append(copied_chapter)

    with tempfile.TemporaryDirectory(prefix="ppt-assets-") as temp_dir_raw:
        temp_dir = Path(temp_dir_raw)
        for position, (title, body) in enumerate(content_sections):
            selected_asset = assets[position] if position < len(assets) else None
            visual_asset = None

            if content_template_index is not None:
                copied_content = _duplicate_slide(prs, content_template_index)
                _clear_slide_text(prs.slides[copied_content])
                if selected_asset is not None:
                    visual_asset = _prepare_visual_asset_for_ppt(selected_asset, temp_dir)
                _populate_template_slide(
                    prs.slides[copied_content],
                    content_template_index,
                    title,
                    body,
                    visual_asset,
                )
                front_indices.append(copied_content)
                continue

            _add_content_slide(prs, title, body)
            front_indices.append(len(prs.slides) - 1)

    if team_template_index is not None:
        copied_team = _duplicate_slide(prs, team_template_index)
        front_indices.append(copied_team)

    if thank_you_template_index is not None:
        copied_thank_you = _duplicate_slide(prs, thank_you_template_index)
        _clear_slide_text(prs.slides[copied_thank_you])
        closing_body = (
            "Next step: confirm owners, close blockers, and approve the next delivery checkpoint status."
        )
        _populate_template_slide(
            prs.slides[copied_thank_you],
            thank_you_template_index,
            "Thank You and Next Decision",
            closing_body,
            None,
        )
        front_indices.append(copied_thank_you)

    _add_backup_divider_slide(prs)
    backup_index = len(prs.slides) - 1

    remaining_indices = [
        idx
        for idx in range(len(prs.slides))
        if idx not in set(front_indices + [backup_index] + template_indices)
    ]
    final_order = front_indices + [backup_index] + template_indices + remaining_indices
    _reorder_slides(prs, final_order)


def _is_stage_markdown(source: Path) -> bool:
    """Return True when source looks like a canonical stage markdown document."""
    if source.suffix.lower() != ".md" or not source.is_file():
        return False
    text = source.read_text(encoding="utf-8", errors="replace")
    required = ("## Vision", "## Goals", "## Constraints")
    return all(token in text for token in required)


def parse_args() -> argparse.Namespace:
    """TODO: add docstring for parse_args."""
    parser = argparse.ArgumentParser(description="Build deterministic deck from source")
    parser.add_argument("--source", required=True)
    parser.add_argument("--repo-root", default=".")
    parser.add_argument("--layer", required=True)
    parser.add_argument("--output")
    return parser.parse_args()


def main() -> int:
    """TODO: add docstring for main."""
    args = parse_args()
    repo_root = Path(args.repo_root).resolve()
    source = Path(args.source).resolve()
    if not source.exists():
        raise SystemExit(f"source not found: {source}")

    template = (
        repo_root
        / ".github/skills/powerpoint/templates"
        / f"{args.layer}_template.pptx"
    )
    template_path, _ = ensure_template(repo_root, args.layer, template)
    out_default = (
        repo_root / "docs/powerpoints" / f"{args.layer}_{source_slug(source)}.pptx"
    )
    output = Path(args.output).resolve() if args.output else out_default

    prs = Presentation(str(template_path))
    is_stage_doc = _is_stage_markdown(source)
    if is_stage_doc:
        sections = _stage_sections(source, args.layer)
    else:
        sections = _collect_chunks(source)
    visuals = _discover_visual_assets(repo_root, source)

    _apply_sections_to_template(prs, sections, visuals)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(
        json.dumps(
            {"status": "ok", "output": str(output), "template": str(template_path)}
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
