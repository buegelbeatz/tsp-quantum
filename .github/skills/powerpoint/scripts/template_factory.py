"""Template creation helpers for layer-specific PowerPoint themes.

Generates a 7-slide standard corporate template:
    1. Title         - dark background, white title + teal subtitle
    2. Agenda        - dark background, 3 standard sections
    3. Chapter Title - dark background, large teal chapter number + white heading
    4. Slide         - light background, dark text placeholder
    5. Our Team      - dark background, all non-generic agents with round portraits
    6. Thank You     - dark background, white closing message + teal tagline
    7. Schema        - light background, colour swatches + typography specification
"""

from __future__ import annotations

from pathlib import Path
from tempfile import NamedTemporaryFile
import math
import random
import shutil
import sys

if str(Path(__file__).resolve().parent) not in sys.path:
    sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation  # type: ignore[import-not-found]
from pptx.dml.color import RGBColor  # type: ignore[import-not-found]
from pptx.enum.text import PP_ALIGN  # type: ignore[import-not-found]
from pptx.util import Inches, Pt  # type: ignore[import-not-found]

from powerpoint_seed import build_seed  # type: ignore[import-not-found]
from powerpoint_svg import make_background_svg  # type: ignore[import-not-found]
from powerpoint_portraits import (  # type: ignore[import-not-found]
    get_portrait_by_index,
)

SLIDE_W = 1280
SLIDE_H = 720

# Text colours for dark vs. light slide variants
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_TEAL = RGBColor(0x16, 0xE0, 0xBD)
_NAVY = RGBColor(0x1F, 0x4E, 0x79)
_DARK = RGBColor(0x07, 0x14, 0x26)


def _format_agent_title(slug: str) -> str:
    """Convert an agent slug to a human-friendly title."""
    acronyms = {
        "ai": "AI",
        "ux": "UX",
        "mcp": "MCP",
        "k3s": "K3S",
    }
    normalized = slug.replace("pullrequest", "pull-request")
    parts = [p for p in normalized.split("-") if p]
    pretty_parts = []
    for part in parts:
        if part in acronyms:
            pretty_parts.append(acronyms[part])
        else:
            pretty_parts.append(part.capitalize())
    return " ".join(pretty_parts)


def _load_non_generic_agent_titles(repo_root: Path) -> list[str]:
    """Load all non-generic agent titles from .github/agents/*.agent.md."""
    agents_dir = repo_root / ".github" / "agents"
    if not agents_dir.exists():
        return []

    titles: list[str] = []
    for path in sorted(agents_dir.glob("*.agent.md")):
        lines = path.read_text(encoding="utf-8").splitlines()
        if not lines or lines[0].strip() != "---":
            continue
        name_value = ""
        for line in lines[1:]:
            if line.strip() == "---":
                break
            if line.startswith("name:"):
                name_value = line.split(":", 1)[1].strip().strip('"').strip("'")
                break
        if not name_value or name_value.startswith("generic-"):
            continue
        titles.append(_format_agent_title(name_value))
    return titles


def _svg_to_png_file(svg: str) -> Path:
    """Convert SVG string to PNG file.

    Uses cairosvg to render SVG markup to a PNG image file.

    Args:
        svg: SVG markup as string.

    Returns:
        Path: Temporary PNG file path.

    Note:
        The caller is responsible for deleting the temporary file
        after use via Path.unlink(missing_ok=True).
    """
    import cairosvg  # type: ignore[import-untyped]

    with NamedTemporaryFile(prefix="ppt-bg-", suffix=".png", delete=False) as tmp:
        cairosvg.svg2png(bytestring=svg.encode("utf-8"), write_to=tmp.name)
        return Path(tmp.name)


def _add_background(slide, png_file: Path) -> None:
    """Add a background image to a slide.

    Args:
        slide: PPTX slide object.
        png_file: Path to PNG file to use as background (fills entire slide).
    """
    slide.shapes.add_picture(
        str(png_file), 0, 0, width=Inches(13.333), height=Inches(7.5)
    )


def _text(
    frame,
    text: str,
    size_pt: int,
    bold: bool,
    color: RGBColor,
    *,
    font_name: str = "Aptos",
    align: PP_ALIGN | None = None,
) -> None:
    """Set text, font size, bold flag, and colour on a text frame.

    Args:
        frame: PPTX text frame object.
        text: Text content to display.
        size_pt: Font size in points.
        bold: Whether to apply bold formatting.
        color: RGB font colour.
    """
    frame.text = text
    run = frame.paragraphs[0].runs[0]
    if align is not None:
        frame.paragraphs[0].alignment = align
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color


def _polaroid_portrait_to_png_file(
    portrait,
    output_size: int = 320,
    corner_radius: int = 24,
    border_px: int = 8,
    inner_padding_px: int = 6,
    trim_ratio: float = 0.06,
) -> Path:
    """Convert portrait into a rounded-rectangle framed PNG with transparent corners.

    Keeps aspect ratio (no stretching), applies one uniform thin white frame, and
    trims outer white frame from the source sprite.
    """
    from PIL import Image, ImageDraw  # type: ignore[import-not-found]

    if portrait.mode != "RGBA":
        portrait = portrait.convert("RGBA")

    # Remove the outer white frame from the sprite portrait before circular masking.
    trim = int(min(portrait.width, portrait.height) * trim_ratio)
    if trim > 0 and portrait.width > 2 * trim and portrait.height > 2 * trim:
        portrait = portrait.crop(
            (trim, trim, portrait.width - trim, portrait.height - trim)
        )

    content_size = max(1, output_size - 2 * border_px)
    image_box_size = max(1, content_size - 2 * inner_padding_px)
    resampling = getattr(Image, "Resampling", Image).LANCZOS
    # Preserve aspect ratio explicitly to avoid any stretching artifacts.
    fitted = portrait.copy()
    fitted.thumbnail((image_box_size, image_box_size), resampling)

    avatar = Image.new("RGBA", (output_size, output_size), (0, 0, 0, 0))
    # Solid white backing avoids dark inner fringes on dark slide backgrounds.
    base_draw = ImageDraw.Draw(avatar)
    base_draw.rounded_rectangle(
        (0, 0, output_size - 1, output_size - 1),
        radius=corner_radius,
        fill=(255, 255, 255, 255),
    )

    image_layer = Image.new(
        "RGBA", (image_box_size, image_box_size), (255, 255, 255, 255)
    )
    img_x = (image_box_size - fitted.width) // 2
    img_y = (image_box_size - fitted.height) // 2 + int(image_box_size * 0.02)
    if img_y + fitted.height > image_box_size:
        img_y = image_box_size - fitted.height
    image_layer.paste(fitted, (img_x, img_y), fitted)

    paste_x = border_px + inner_padding_px
    paste_y = border_px + inner_padding_px
    avatar.paste(image_layer, (paste_x, paste_y), image_layer)

    mask = Image.new("L", (output_size, output_size), 0)
    draw = ImageDraw.Draw(mask)
    draw.rounded_rectangle(
        (0, 0, output_size - 1, output_size - 1), radius=corner_radius, fill=255
    )
    avatar.putalpha(mask)

    # Single thin, uniform white border around the entire image frame.
    border_draw = ImageDraw.Draw(avatar)
    border_draw.rounded_rectangle(
        (
            border_px // 2,
            border_px // 2,
            output_size - 1 - border_px // 2,
            output_size - 1 - border_px // 2,
        ),
        radius=max(1, corner_radius - border_px // 2),
        outline=(255, 255, 255, 255),
        width=border_px,
    )

    with NamedTemporaryFile(
        prefix="ppt-polaroid-portrait-", suffix=".png", delete=False
    ) as tmp:
        avatar.save(tmp.name, format="PNG")
        return Path(tmp.name)


def _add_title_slide(slide, layer: str, repo_name: str) -> None:
    """Slide 1 – Title: large white layer name + teal repository subtitle.

    Args:
        slide: PPTX slide object (dark background expected).
        layer: Layer display name shown as the main title.
        repo_name: Repository name shown as subtitle.
    """
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.6))
    _text(box.text_frame, layer, 54, True, _WHITE)

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.2), Inches(0.8))
    _text(sub.text_frame, f"Repository: {repo_name}", 22, False, _TEAL)


def _add_agenda_slide(slide) -> None:
    """Slide 2 – Agenda: dark mode with 3 standard sections."""
    heading = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.9)
    )
    _text(heading.text_frame, "Agenda", 40, True, _WHITE)

    sections = [
        "1. Context and Objectives",
        "2. Solution Overview",
        "3. Next Steps and Decisions",
    ]
    start_y = 2.1
    for idx, section in enumerate(sections):
        row_y = start_y + idx * 1.2
        bullet = slide.shapes.add_textbox(
            Inches(1.0), Inches(row_y), Inches(0.5), Inches(0.7)
        )
        _text(bullet.text_frame, "•", 30, True, _TEAL)

        item = slide.shapes.add_textbox(
            Inches(1.5), Inches(row_y + 0.05), Inches(10.7), Inches(0.7)
        )
        _text(item.text_frame, section, 27, True, _WHITE)


def _add_content_slide(slide) -> None:
    """Slide 4 – Slide: dark-navy section title + dark body placeholder.

    Args:
        slide: PPTX slide object (light background expected).
    """
    title = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.9)
    )
    _text(title.text_frame, "Section Title", 32, True, _NAVY)

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.6), Inches(5.3))
    _text(body.text_frame, "Replace this with your slide content.", 18, False, _DARK)


def _add_chapter_slide(slide) -> None:
    """Slide 3 – Chapter Title: large teal chapter number + white chapter heading.

    Args:
        slide: PPTX slide object (dark background expected).
    """
    num = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(2.0))
    _text(num.text_frame, "01", 96, True, _TEAL)

    heading = slide.shapes.add_textbox(
        Inches(0.8), Inches(3.8), Inches(11.5), Inches(1.2)
    )
    _text(heading.text_frame, "Chapter Title", 40, True, _WHITE)


def _add_team_slide(slide, seed: int, repo_root: Path) -> None:
    """Slide 5 – Our Team: all non-generic agents in polaroid-style cards.

    Args:
        slide: PPTX slide object (dark background expected).
        seed: Base seed for deterministic portrait selection.
    """
    title = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.35), Inches(11.5), Inches(0.8)
    )
    _text(title.text_frame, "Our Team", 32, True, _WHITE)

    agent_titles = _load_non_generic_agent_titles(repo_root)
    if not agent_titles:
        agent_titles = ["Agile Coach", "Fullstack Engineer", "AI Expert"]

    count = len(agent_titles)
    cols = 5 if count >= 13 else 4 if count >= 9 else 3
    rows = math.ceil(count / cols)

    area_x = 0.6
    area_y = 1.25
    area_w = 12.1
    area_h = 5.95

    gap_x = 0.28
    cell_w = (area_w - gap_x * (cols - 1)) / cols
    row_h = area_h / rows
    frame_size = min(1.42, cell_w * 0.74, row_h * 0.68)
    label_h = 0.3
    portrait_indices = _unique_portrait_indices(seed + 113, count)

    for i, agent_title in enumerate(agent_titles):
        row = i // cols
        col = i % cols
        cell_x = area_x + col * (cell_w + gap_x)
        cell_y = area_y + row * row_h
        frame_x = cell_x + (cell_w - frame_size) / 2
        frame_y = cell_y + max(0.01, (row_h - (frame_size + label_h + 0.08)) / 2)

        portrait = get_portrait_by_index(portrait_indices[i])
        portrait_file = _polaroid_portrait_to_png_file(portrait)
        try:
            slide.shapes.add_picture(
                str(portrait_file),
                Inches(frame_x),
                Inches(frame_y),
                width=Inches(frame_size),
                height=Inches(frame_size),
            )
        finally:
            portrait_file.unlink(missing_ok=True)

        label = slide.shapes.add_textbox(
            Inches(cell_x),
            Inches(frame_y + frame_size + 0.08),
            Inches(cell_w),
            Inches(label_h),
        )
        _text(label.text_frame, agent_title, 12, True, _WHITE, align=PP_ALIGN.CENTER)


def _unique_portrait_indices(seed: int, count: int, pool_size: int = 25) -> list[int]:
    """Return deterministic portrait indices with no duplicates until pool exhausted."""
    rng = random.Random(seed)
    pool = list(range(pool_size))
    indices: list[int] = []
    while len(indices) < count:
        shuffled = pool[:]
        rng.shuffle(shuffled)
        indices.extend(shuffled)
    return indices[:count]


def _add_thanx_slide(slide) -> None:
    """Slide 5 – Thank You: centred white closing message + teal tagline.

    Args:
        slide: PPTX slide object (dark background expected).
    """
    msg = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.0), Inches(2.0))
    _text(msg.text_frame, "Thank You", 64, True, _WHITE)

    tag = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(11.0), Inches(0.8))
    _text(tag.text_frame, "Questions? Let's talk.", 24, False, _TEAL)


def _add_color_schema_slide(slide) -> None:
    """Slide 7 – Schema: navy heading + colour swatches + typography specification.

    Swatches are 1.6"×1.6" rectangles with 0.35" gaps, centred horizontally.

    Args:
        slide: PPTX slide object (light background expected).
    """
    heading = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8)
    )
    _text(heading.text_frame, "Schema", 30, True, _NAVY)

    swatches = [
        ("#071426", "Dark Navy"),
        ("#0b2a4a", "Deep Blue"),
        ("#16e0bd", "Teal"),
        ("#59b3ff", "Sky Blue"),
        ("#1f4e79", "Navy"),
        ("#2a9d8f", "Green Teal"),
    ]
    swatch_w, swatch_h = 1.35, 1.35
    gap = 0.35
    num = len(swatches)
    start_x = (13.333 - (num * swatch_w + (num - 1) * gap)) / 2
    start_y = 1.4

    for i, (hex_color, label) in enumerate(swatches):
        x = start_x + i * (swatch_w + gap)
        r = int(hex_color[1:3], 16)
        g = int(hex_color[3:5], 16)
        b = int(hex_color[5:7], 16)
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(x),
            Inches(start_y),
            Inches(swatch_w),
            Inches(swatch_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        shape.line.color.rgb = RGBColor(r, g, b)

        lbl = slide.shapes.add_textbox(
            Inches(x), Inches(start_y + swatch_h + 0.1), Inches(swatch_w), Inches(0.5)
        )
        _text(lbl.text_frame, label, 12, False, _NAVY)

    spec_heading = slide.shapes.add_textbox(
        Inches(0.8), Inches(3.45), Inches(11.5), Inches(0.5)
    )
    _text(
        spec_heading.text_frame,
        "Typography Specification (Live Samples)",
        20,
        True,
        _NAVY,
    )

    sample_1 = slide.shapes.add_textbox(
        Inches(0.95), Inches(3.95), Inches(11.2), Inches(0.7)
    )
    _text(
        sample_1.text_frame,
        "Aptos Display Bold 34 pt  |  Enterprise Presentation Title",
        28,
        True,
        _DARK,
        font_name="Aptos Display",
    )

    sample_2 = slide.shapes.add_textbox(
        Inches(0.95), Inches(4.62), Inches(11.2), Inches(0.62)
    )
    _text(
        sample_2.text_frame,
        "Aptos SemiBold 22 pt  |  Agenda, chapter and key section headings",
        20,
        True,
        _NAVY,
        font_name="Aptos",
    )

    sample_3 = slide.shapes.add_textbox(
        Inches(0.95), Inches(5.22), Inches(11.2), Inches(0.62)
    )
    _text(
        sample_3.text_frame,
        "Aptos Regular 16 pt  |  Body text and explanatory content",
        16,
        False,
        _DARK,
        font_name="Aptos",
    )

    sample_4 = slide.shapes.add_textbox(
        Inches(0.95), Inches(5.8), Inches(11.2), Inches(0.56)
    )
    _text(
        sample_4.text_frame,
        "Fallback for compatibility: Calibri (Office default fallback)",
        13,
        False,
        _NAVY,
        font_name="Calibri",
    )


def _is_background_picture(shape) -> bool:
    """Return True when shape is a full-slide background picture positioned at (0, 0)."""
    try:
        return (
            shape.left == 0
            and shape.top == 0
            and abs(shape.width - Inches(13.333)) < 5_000
            and abs(shape.height - Inches(7.5)) < 5_000
        )
    except AttributeError:
        return False


def _clear_content_shapes(slide) -> None:
    """Remove every shape from a slide except the full-slide background picture."""
    sp_tree = slide.shapes._spTree
    for elem in [s._element for s in list(slide.shapes) if not _is_background_picture(s)]:
        sp_tree.remove(elem)


def _create_template_from_scratch(
    repo_root: Path,
    layer: str,
    template_path: Path,
    seed: int,
    repo_name: str,
) -> Path:
    """Generate all 7 slides from a blank Presentation when no base template exists."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 7-slide spec: (SVG variant, seed offset)
    slides_spec = [
        ("title", 0),  # 1 Title
        ("title", 7),  # 2 Agenda
        ("title", 14),  # 3 Chapter Title
        ("content", 21),  # 4 Slide
        ("title", 28),  # 5 Our Team
        ("title", 35),  # 6 Thank You
        ("content", 42),  # 7 Schema
    ]
    layout = prs.slide_layouts[6]  # blank
    created_slides = []
    for variant, offset in slides_spec:
        slide = prs.slides.add_slide(layout)
        bg_png = _svg_to_png_file(
            make_background_svg(SLIDE_W, SLIDE_H, seed + offset, variant)
        )
        try:
            _add_background(slide, bg_png)
        finally:
            bg_png.unlink(missing_ok=True)
        created_slides.append(slide)

    _add_title_slide(created_slides[0], layer, repo_name)
    _add_agenda_slide(created_slides[1])
    _add_chapter_slide(created_slides[2])
    _add_content_slide(created_slides[3])
    _add_team_slide(created_slides[4], seed, repo_root)
    _add_thanx_slide(created_slides[5])
    _add_color_schema_slide(created_slides[6])

    template_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(template_path))
    return template_path


def create_template(repo_root: Path, layer: str, template_path: Path) -> Path:
    """Clone the base template and customise it for the requested layer.

    Copies *digital-generic-team_template.pptx* (the hand-crafted master) and
    refreshes only the two layer-specific slides (Title, Our Team) so that all
    background graphics, colour schemes, and typography from the master are
    preserved.  Falls back to from-scratch generation when the master template
    is absent or when the target path equals the source.

    Args:
        repo_root: Repository root directory for seed and agent discovery.
        layer: Layer identifier used for the title slide subtitle.
        template_path: Output path for the resulting template file.

    Returns:
        Path: Path to the created template file.
    """
    template_path.parent.mkdir(parents=True, exist_ok=True)
    repo_name = repo_root.name
    seed = build_seed(repo_name, layer)

    script_dir = Path(__file__).resolve().parent
    base_template = script_dir.parent / "templates" / "digital-generic-team_template.pptx"

    if base_template.exists() and base_template.resolve() != template_path.resolve():
        shutil.copy2(str(base_template), str(template_path))
        prs = Presentation(str(template_path))
        if len(prs.slides) > 0:
            _clear_content_shapes(prs.slides[0])
            _add_title_slide(prs.slides[0], layer, repo_name)
        if len(prs.slides) > 4:
            _clear_content_shapes(prs.slides[4])
            _add_team_slide(prs.slides[4], seed, repo_root)
        prs.save(str(template_path))
        return template_path

    return _create_template_from_scratch(repo_root, layer, template_path, seed, repo_name)


def ensure_template(
    repo_root: Path, layer: str, template_path: Path
) -> tuple[Path, bool]:
    """Return existing template or create a new one if missing.

    Args:
        repo_root: Repository root directory.
        layer: Layer name for the template.
        template_path: Path where template should be created.

    Returns:
        tuple: (template_path, created) where created is True if template was newly created.
    """
    if template_path.exists():
        return template_path, False
    return create_template(repo_root, layer, template_path), True
